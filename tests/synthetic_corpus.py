"""A fictional Estonian legal corpus, generated rather than committed.

Every byte the extraction tests read is produced here, at test time, from code
in this file. Nothing is a checked-in binary, which buys three things:

* **Nothing real can leak in.** A generated corpus cannot accidentally become a
  real ministry PDF somebody dropped into `tests/fixtures/` while debugging.
  The repository is public and the rule is absolute (Stage-2B brief 60, 84).
* **The fixtures are readable.** "What is on page 4" is answered by reading the
  function, not by opening a binary in a hex editor.
* **Each test can ask for the shape it needs** — a corrupt PDF, a scanned page,
  a workbook with a formula — instead of the corpus growing a file per case.

The PDF writer is hand-rolled rather than pulled from a library. It is fifty
lines, it produces exactly the bytes the tests reason about, and it keeps the
dependency list to parsers we actually ship. `Pillow` writes the image-only PDF
because rasterising is what the scanned case *is*.

Every name, title and address below is invented. `example.invalid` and
`.invalid` domains cannot resolve, by RFC.
"""

from __future__ import annotations

import io
import struct
import zipfile
from dataclasses import dataclass
from email.message import EmailMessage

# -- vocabulary ------------------------------------------------------------
#
# Deliberately plausible and deliberately fictional. "Näidisministeerium" is not
# a ministry; the draft it is consulting on does not exist.

MINISTRY = "Näidisministeerium"
DRAFT_TITLE = "Pakendiaruandluse katse-eelnõu"

#: Words placed in exactly one place each, so a search test can prove *which*
#: source answered rather than merely that something did (Stage-2B brief 65).
ONLY_ON_PDF_PAGE_4 = "kaubaaluste"
ONLY_IN_DOCX_TABLE = "vahearuandlus"
ONLY_ON_XLSX_SHEET_2 = "hoiustamiskulu"
ONLY_ON_PPTX_SLIDE_3 = "üleminekutähtaeg"
ONLY_IN_EMAIL_BODY = "kooskõlastusringile"
ONLY_IN_ATTACHMENT = "pandipakend"
ONLY_IN_OCR_IMAGE = "MAARUS"


# -- PDF -------------------------------------------------------------------


def _pdf(objects: list[bytes], *, trailer_extra: str = "") -> bytes:
    """Assemble numbered objects into a valid PDF with a correct xref table.

    The cross-reference offsets are computed from the bytes actually written,
    which is the only part of this that has to be exact: a PDF with a wrong
    xref is precisely the "corrupt file" case, and producing one by accident
    while trying to produce a healthy one would make two tests lie at once.
    """
    out = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{number} 0 obj\n".encode() + body + b"\nendobj\n"
    start = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets[1:]:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<</Size {len(objects) + 1}/Root 1 0 R{trailer_extra}>>\n"
        f"startxref\n{start}\n%%EOF\n"
    ).encode()
    return bytes(out)


def text_pdf(pages: list[str]) -> bytes:
    """A text-native PDF, one content stream per page.

    WinAnsi encoding with a standard Type 1 font, so `õ` and `ä` survive the
    round trip. An Estonian corpus that tested only ASCII would prove nothing
    about the corpus we actually have.
    """
    page_count = len(pages)
    first_page_object = 4
    kids = " ".join(f"{first_page_object + index * 2} 0 R" for index in range(page_count))

    objects: list[bytes] = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        f"<</Type/Pages/Kids[{kids}]/Count {page_count}>>".encode(),
        b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica/Encoding/WinAnsiEncoding>>",
    ]
    for index, body in enumerate(pages):
        content_object = first_page_object + index * 2 + 1
        objects.append(
            (
                "<</Type/Page/Parent 2 0 R/MediaBox[0 0 595 842]"
                "/Resources<</Font<</F1 3 0 R>>>>"
                f"/Contents {content_object} 0 R>>"
            ).encode()
        )
        lines = body.splitlines() or [""]
        drawn = "\n".join(f"({_escape(line)}) Tj 0 -16 Td" for line in lines)
        stream = f"BT /F1 12 Tf 50 780 Td\n{drawn}\nET".encode("cp1252", errors="replace")
        objects.append(f"<</Length {len(stream)}>>\nstream\n".encode() + stream + b"\nendstream")
    return _pdf(objects)


def _escape(line: str) -> str:
    return line.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")


def government_pdf() -> bytes:
    """Six pages of a fictional consultation letter. Page 4 is the marked one."""
    return text_pdf(
        [
            f"{MINISTRY}\nKooskõlastuskiri\n\n{DRAFT_TITLE}",
            "1. Eelnõu eesmärk\n\nEelnõu korrastab pakendiaruandluse tähtaegu.",
            "2. Mõju ettevõtjatele\n\nAruandluskoormus väheneb väikeettevõtjatel.",
            f"3. Erisused\n\nErisus kohaldub {ONLY_ON_PDF_PAGE_4} ringlusele.",
            "4. Rakendamine\n\nSäte jõustub üldises korras.",
            "5. Kooskõlastamine\n\nPalume seisukohta kolme nädala jooksul.",
        ]
    )


def scanned_pdf(text: str = ONLY_IN_OCR_IMAGE) -> bytes:
    """An image-only PDF: a page with no text layer at all.

    This is what a photographed annex looks like to a parser, and it is the
    only honest way to test the OCR trigger — a PDF whose text layer is merely
    short would still take the native path.
    """

    image = _text_image(text)
    buffer = io.BytesIO()
    image.save(buffer, format="PDF", resolution=200.0)
    return buffer.getvalue()


def mixed_pdf() -> bytes:
    """Not built as one file, because a PDF cannot hold both cheaply here.

    The mixed case — a typed cover page and a photographed annex in one file —
    is exercised by the PDF parser's per-page decision, which is tested against
    a text PDF whose later pages are deliberately empty.
    """
    return text_pdf(
        [
            f"{MINISTRY}\nKaaskiri\n\n{DRAFT_TITLE}",
            "",
            "",
        ]
    )


def corrupt_pdf() -> bytes:
    """A file that claims to be a PDF and is not.

    The header is right — so the upload allowlist accepts it, which is the
    point — and everything after it is noise.
    """
    return b"%PDF-1.4\n" + b"\x00\xff" * 200


def encrypted_pdf() -> bytes:
    """A PDF with an encryption dictionary and no usable key.

    Enough to make `is_encrypted` true and decryption fail, which is the branch
    under test. Nothing here attempts real RC4 or AES: the parser must refuse
    without guessing a password, so a genuinely encrypted file would be testing
    somebody else's cryptography.
    """
    objects: list[bytes] = [
        b"<</Type/Catalog/Pages 2 0 R>>",
        b"<</Type/Pages/Kids[3 0 R]/Count 1>>",
        b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 595 842]/Contents 4 0 R>>",
        b"<</Length 0>>\nstream\n\nendstream",
        # A standard security handler whose owner and user keys are zeroes, so
        # the empty-password path the parser tries is rejected and it has to
        # give up rather than guess.
        (f"<</Filter/Standard/V 1/R 2/P -1/O <{'0' * 64}>/U <{'0' * 64}>>>").encode(),
    ]
    return _pdf(objects, trailer_extra="/Encrypt 5 0 R/ID[<00><00>]")


# -- Office ----------------------------------------------------------------


def draft_docx() -> bytes:
    from docx import Document

    document = Document()
    document.add_heading(DRAFT_TITLE, level=1)
    document.add_paragraph("Koja märkused eelnõu kohta.")
    document.add_paragraph("Ettepanek on tähtaega pikendada.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Säte"
    table.cell(0, 1).text = "Märkus"
    table.cell(1, 0).text = "§ 12"
    table.cell(1, 1).text = f"Ettepanek: {ONLY_IN_DOCX_TABLE} kaotada"
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def annex_xlsx() -> bytes:
    import openpyxl

    workbook = openpyxl.Workbook()
    first = workbook.active
    first.title = "Kokkuvõte"
    first["A1"] = "Näitaja"
    first["B1"] = "Väärtus"
    first["A2"] = "Ettevõtjaid"
    first["B2"] = 412
    # A formula with no cached value. openpyxl in data_only mode returns None
    # for it, which is exactly the behaviour the parser documents: the workbook
    # is read as data and nothing is evaluated.
    first["B3"] = "=B2*2"

    second = workbook.create_sheet("Kulud")
    second["A1"] = "Kuluartikkel"
    second["A2"] = ONLY_ON_XLSX_SHEET_2
    second["B2"] = 1840

    workbook.create_sheet("Tühi")
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def briefing_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for index, text in enumerate(
        [
            f"{DRAFT_TITLE} — ülevaade",
            "Koja seisukoht: toetame põhimõtet",
            f"Tähelepanu: {ONLY_ON_PPTX_SLIDE_3}",
        ],
        start=1,
    ):
        slide = presentation.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        box.text_frame.text = text
        if index == 2:
            slide.notes_slide.notes_text_frame.text = "Esineja märkus: mitte lubada erandit."
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def zip_archive() -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("lisa.txt", "Näidissisu")
    return buffer.getvalue()


def zip_bomb(members: int = 5000) -> bytes:
    """Many members, each trivially compressible. Trips the member ceiling."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        for index in range(members):
            archive.writestr(f"f{index}.txt", "0" * 2000)
    return buffer.getvalue()


def broken_docx() -> bytes:
    """A ZIP that is a valid container and not an Office document."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("hello.txt", "not a document")
    return buffer.getvalue()


# -- text ------------------------------------------------------------------


def memo_txt(encoding: str = "utf-8") -> bytes:
    return (
        "Märkmed kohtumiselt\n\n"
        "Ministeerium lubas tähtaega pikendada.\n"
        "Järgmine samm: koostada Koja arvamus.\n"
    ).encode(encoding)


def undecodable_txt() -> bytes:
    """Bytes that are not valid in any encoding the parser tries."""
    return bytes([0x80, 0x81, 0xFE, 0xFF] * 64)


def counts_csv() -> bytes:
    return ("Ettevõtja;Pakendiliik;Kogus\nNäidis AS;plast;120\nNäidis OÜ;klaas;45\n").encode()


# -- images ----------------------------------------------------------------


def _text_image(text: str, width: int = 900, height: int = 260):
    """A white image with large black text, legible to Tesseract.

    The default PIL font is bitmap and small, so it is scaled up rather than
    relying on a TrueType file being present — a fixture that needs a system
    font is a fixture that fails on somebody else's machine.
    """
    from PIL import Image, ImageDraw

    small = Image.new("RGB", (width // 4, height // 4), "white")
    draw = ImageDraw.Draw(small)
    draw.text((6, height // 12), text, fill="black")
    return small.resize((width, height), Image.LANCZOS)


def scanned_png(text: str = ONLY_IN_OCR_IMAGE) -> bytes:
    buffer = io.BytesIO()
    _text_image(text).save(buffer, format="PNG")
    return buffer.getvalue()


def photo_jpeg() -> bytes:
    from PIL import Image

    buffer = io.BytesIO()
    Image.new("RGB", (240, 160), (90, 110, 130)).save(buffer, format="JPEG")
    return buffer.getvalue()


# -- email -----------------------------------------------------------------


@dataclass(frozen=True)
class EmailFixture:
    content: bytes
    subject: str
    sender: str


def consultation_eml(
    *, with_html: bool = False, attachments: bool = True, inline_logo: bool = True
) -> bytes:
    """A ministry consultation message, in the shape they actually arrive.

    Covering text, a PDF annex, and — when asked for — a signature logo as an
    inline resource, because separating those two is a rule the parser has to
    get right and a fixture without one cannot test.
    """
    message = EmailMessage()
    message["Subject"] = f"{DRAFT_TITLE} — kooskõlastamiseks"
    message["From"] = "Kadri Näidis <kadri@naidisministeerium.invalid>"
    message["To"] = "Koja õigusosakond <oigus@koda.invalid>"
    message["Cc"] = "arhiiv@koda.invalid"
    message["Date"] = "Mon, 17 Aug 2026 09:12:00 +0300"
    message["Message-ID"] = "<katse-1@naidisministeerium.invalid>"
    message["In-Reply-To"] = "<katse-0@naidisministeerium.invalid>"
    message["References"] = "<katse-0@naidisministeerium.invalid>"

    body = f"Tere\n\nSaadame {DRAFT_TITLE} {ONLY_IN_EMAIL_BODY}.\n\nLugupidamisega\nKadri Näidis\n"
    message.set_content(body)
    if with_html:
        # The script tag is the point of this branch: the sanitiser has to
        # remove the element and its contents, not merely strip the tags.
        message.add_alternative(
            "<html><head><style>p{color:red}</style></head><body>"
            "<script>alert('xss')</script>"
            f"<p>Saadame {DRAFT_TITLE} {ONLY_IN_EMAIL_BODY}.</p>"
            '<img src="https://tracker.invalid/pixel.gif">'
            "</body></html>",
            subtype="html",
        )

    if attachments:
        message.add_attachment(
            text_pdf([f"Lisa 1\n\nMääratlus: {ONLY_IN_ATTACHMENT} on tagatisrahaga pakend."]),
            maintype="application",
            subtype="pdf",
            filename="lisa-1.pdf",
        )
        message.add_attachment(
            memo_txt(),
            maintype="text",
            subtype="plain",
            filename="markmed.txt",
        )
    if inline_logo:
        part = message.add_attachment(
            scanned_png("LOGO"),
            maintype="image",
            subtype="png",
            filename="allkiri-logo.png",
            cid="<logo@naidisministeerium.invalid>",
        )
        del part  # add_attachment returns None; the cid is what matters
    return message.as_bytes()


def malformed_eml() -> bytes:
    """Headers that stop mid-way and a body that never starts."""
    return "Subject: katkine\r\nFrom: keegi\r\n\x00\x00binaarne prügi".encode()


def outlook_msg(*, attachment: bool = True) -> bytes:
    """A synthetic Outlook message, built with `extract-msg`'s own OLE writer.

    Writing one rather than committing one keeps the corpus generated end to
    end. The property streams below are the minimum a real `.msg` carries for
    the fields Juristid extracts; anything more would be testing Outlook rather
    than this parser.
    """
    from extract_msg.ole_writer import OleWriter

    def unicode_property(text: str) -> bytes:
        return text.encode("utf-16-le")

    streams: dict[str, bytes] = {
        # PR_MESSAGE_CLASS — without it the file is an OLE container and not a
        # message, and `openMsg` says so.
        "__substg1.0_001A001F": unicode_property("IPM.Note"),
        "__substg1.0_0037001F": unicode_property(f"{DRAFT_TITLE} — Outlooki koopia"),
        "__substg1.0_1000001F": unicode_property(
            f"Tere\r\n\r\nSaadame {DRAFT_TITLE} {ONLY_IN_EMAIL_BODY}.\r\n"
        ),
        "__substg1.0_0C1A001F": unicode_property("Kadri Näidis"),
        "__substg1.0_0E04001F": unicode_property("Koja õigusosakond"),
        "__substg1.0_007D001F": unicode_property(
            "Message-ID: <katse-2@naidisministeerium.invalid>\r\n"
            "Date: Mon, 17 Aug 2026 09:12:00 +0300\r\n"
        ),
    }

    writer = OleWriter()
    header = b"\x00" * 8 + struct.pack("<IIII", 0, 0, 0, 1 if attachment else 0) + b"\x00" * 8
    properties = bytearray(header)
    for name, payload in streams.items():
        tag = int(name.rsplit("_", 1)[-1], 16)
        properties += struct.pack("<IIII", tag, 6, len(payload) + 2, 0)
    writer.addEntry("__properties_version1.0", bytes(properties))
    for name, payload in streams.items():
        writer.addEntry(name, payload)

    # An empty named-property map. Real messages always have one and some
    # readers assume it.
    writer.addEntry("__nameid_version1.0", storage=True)
    for tag in ("00020102", "00030102", "00040102"):
        writer.addEntry(f"__nameid_version1.0/__substg1.0_{tag}", b"")

    if attachment:
        payload = text_pdf([f"Lisa 1\n\nMääratlus: {ONLY_IN_ATTACHMENT} on tagatisrahaga pakend."])
        folder = "__attach_version1.0_#00000000"
        attachment_streams = {
            f"{folder}/__substg1.0_3707001F": unicode_property("lisa-1.pdf"),
            f"{folder}/__substg1.0_3704001F": unicode_property("lisa-1"),
            f"{folder}/__substg1.0_370E001F": unicode_property("application/pdf"),
            f"{folder}/__substg1.0_37010102": payload,
        }
        attachment_properties = bytearray(b"\x00" * 8)
        for name, blob in attachment_streams.items():
            tag = int(name.rsplit("_", 1)[-1], 16)
            length = len(blob) + (2 if name.endswith("001F") else 0)
            attachment_properties += struct.pack("<IIII", tag, 6, length, 0)
        writer.addEntry(folder, storage=True)
        writer.addEntry(f"{folder}/__properties_version1.0", bytes(attachment_properties))
        for name, blob in attachment_streams.items():
            writer.addEntry(name, blob)

    buffer = io.BytesIO()
    writer.write(buffer)
    return buffer.getvalue()
