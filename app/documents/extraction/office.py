"""DOCX, XLSX and PPTX: read as data, never opened as programs.

All three are ZIP containers full of XML, and all three are read here with
libraries that parse that XML directly. Nothing invokes Word, Excel,
PowerPoint or LibreOffice; nothing evaluates a formula; nothing resolves an
embedded object or a linked image. A document that wants to run something on
open gets exactly the same treatment as one that does not, which is the only
version of this that is defensible (Stage-2B brief 15).

Locators differ per format because the formats genuinely know different things:

* PPTX knows its slide numbers, so a match reports `slaid 12`.
* XLSX knows its sheet names and cell references, so a match reports the sheet
  and the range.
* DOCX knows **neither**. Word paginates at render time against a specific
  printer and font set; the file does not contain page boundaries. So a DOCX
  match reports a section, a paragraph range or a table, and never a page. An
  invented `lk 4` sends a lawyer to the wrong place in a 90-page draft, and one
  such experience costs more trust than the locator ever earns
  (Stage-2B brief 16).
"""

from __future__ import annotations

import io

from app.documents.enums import LocatorKind
from app.documents.extraction.base import (
    DerivativePayload,
    Fragment,
    ParseResult,
    SourceFile,
    normalise,
    registry,
)
from app.documents.extraction.errors import ExtractionFailed
from app.documents.extraction.limits import (
    current_limits,
    guard_character_budget,
    guard_office_container,
)

#: How many paragraphs share one DOCX fragment. Small enough that a locator is
#: specific, large enough that a 400-paragraph draft is not 400 index rows.
DOCX_PARAGRAPHS_PER_FRAGMENT = 40


class DocxParser:
    name = "docx"
    version = "1"
    mime_types = frozenset(
        {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
    )

    def parse(self, source: SourceFile) -> ParseResult:
        import docx

        limits = current_limits()
        handle = io.BytesIO(source.content)
        guard_office_container(handle, limits)

        try:
            document = docx.Document(handle)
        except Exception as error:
            raise ExtractionFailed(
                "unreadable_docx", "DOCX-i ei õnnestunud avada; fail võib olla rikutud."
            ) from error

        fragments: list[Fragment] = []
        total = 0
        ordinal = 0

        paragraphs = [normalise(p.text) for p in document.paragraphs]
        first_index = 1
        buffer: list[str] = []
        for index, text in enumerate(paragraphs, start=1):
            if text:
                buffer.append(text)
            if len(buffer) >= DOCX_PARAGRAPHS_PER_FRAGMENT or index == len(paragraphs):
                if buffer:
                    body = "\n".join(buffer)
                    total += len(body)
                    guard_character_budget(total, limits)
                    ordinal += 1
                    fragments.append(
                        Fragment(
                            text=body,
                            locator_kind=LocatorKind.SECTION,
                            locator={"paragraph_from": first_index, "paragraph_to": index},
                            locator_label=f"lõigud {first_index}–{index}",
                        )
                    )
                    buffer = []
                first_index = index + 1

        for number, table in enumerate(document.tables, start=1):
            body = _table_text(table)
            if not body:
                continue
            total += len(body)
            guard_character_budget(total, limits)
            fragments.append(
                Fragment(
                    text=body,
                    locator_kind=LocatorKind.SECTION,
                    locator={"table": number},
                    locator_label=f"tabel {number}",
                )
            )

        # Headers and footers carry the file reference and the ministry's own
        # registration number more often than they carry anything else, and
        # those are exactly what somebody searches for.
        for number, section in enumerate(document.sections, start=1):
            for label, part in (("päis", section.header), ("jalus", section.footer)):
                body = normalise("\n".join(p.text for p in part.paragraphs))
                if not body:
                    continue
                total += len(body)
                guard_character_budget(total, limits)
                fragments.append(
                    Fragment(
                        text=body,
                        locator_kind=LocatorKind.SECTION,
                        locator={"section": number, "part": label},
                        locator_label=f"{label} (osa {number})",
                    )
                )

        if not fragments:
            raise ExtractionFailed("no_text", "DOCX-ist ei leitud teksti.")

        return ParseResult(
            derivatives=(
                DerivativePayload(
                    fragments=tuple(fragments),
                    metadata={
                        "paragraph_count": len(paragraphs),
                        "table_count": len(document.tables),
                        "section_count": len(document.sections),
                        # Said explicitly so nobody later reads the absence of a
                        # page locator as a parser that forgot.
                        "pagination": "not-available-in-format",
                    },
                ),
            )
        )


def _table_text(table: object) -> str:
    rows: list[str] = []
    for row in table.rows:  # type: ignore[attr-defined]
        cells = [normalise(cell.text) for cell in row.cells]
        line = " | ".join(cell for cell in cells if cell)
        if line:
            rows.append(line)
    return "\n".join(rows)


class XlsxParser:
    name = "xlsx"
    version = "1"
    mime_types = frozenset({"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"})

    def parse(self, source: SourceFile) -> ParseResult:
        import openpyxl

        limits = current_limits()
        handle = io.BytesIO(source.content)
        guard_office_container(handle, limits)

        try:
            # data_only=True asks for the values Excel last cached, which is a
            # read of stored data. It is emphatically not evaluation: openpyxl
            # has no formula engine, and a cell whose cached value is absent
            # simply has none (Stage-2B brief 17).
            workbook = openpyxl.load_workbook(handle, data_only=True, read_only=True)
        except Exception as error:
            raise ExtractionFailed(
                "unreadable_xlsx", "XLSX-i ei õnnestunud avada; fail võib olla rikutud."
            ) from error

        fragments: list[Fragment] = []
        total = 0
        try:
            for sheet in workbook.worksheets:
                total = self._read_sheet(sheet, fragments, total, limits)
        finally:
            workbook.close()

        if not fragments:
            raise ExtractionFailed("no_text", "Töövihikust ei leitud sisu.")

        return ParseResult(
            derivatives=(
                DerivativePayload(
                    fragments=tuple(fragments),
                    metadata={
                        "sheet_count": len(workbook.sheetnames),
                        "sheets": list(workbook.sheetnames),
                        "formula_values": "cached-only",
                    },
                ),
            )
        )

    def _read_sheet(
        self, sheet: object, fragments: list[Fragment], total: int, limits: object
    ) -> int:
        """One sheet, grouped into bounded row blocks.

        A worksheet whose used range is a million mostly-empty rows is common —
        somebody once pressed Ctrl+Down — and turning that into a million
        fragments would be a self-inflicted denial of service. Blank cells
        contribute nothing, and a block that produced no text produces no
        fragment.
        """
        rows_per_block = 50
        block: list[str] = []
        first_row = 1
        produced = 0
        row_number = 0

        for row_number, row in enumerate(sheet.iter_rows(values_only=True), start=1):  # type: ignore[attr-defined]
            line = " | ".join(_cell_text(value) for value in row if _cell_text(value))
            if line:
                block.append(f"{row_number}: {line}")
            if row_number - first_row + 1 >= rows_per_block:
                total, produced = self._flush(
                    sheet, block, first_row, row_number, fragments, total, limits, produced
                )
                block = []
                first_row = row_number + 1
            if produced >= limits.max_xlsx_fragments_per_sheet:  # type: ignore[attr-defined]
                raise ExtractionFailed(
                    "sheet_fragment_limit",
                    f'Tööleht "{sheet.title}" ületab lubatud osade arvu.',  # type: ignore[attr-defined]
                )

        if block:
            total, produced = self._flush(
                sheet, block, first_row, row_number, fragments, total, limits, produced
            )
        return total

    def _flush(
        self,
        sheet: object,
        block: list[str],
        first_row: int,
        last_row: int,
        fragments: list[Fragment],
        total: int,
        limits: object,
        produced: int,
    ) -> tuple[int, int]:
        if not block:
            return total, produced
        body = "\n".join(block)
        total += len(body)
        guard_character_budget(total, limits)  # type: ignore[arg-type]
        name = sheet.title  # type: ignore[attr-defined]
        fragments.append(
            Fragment(
                text=body,
                locator_kind=LocatorKind.SHEET,
                locator={"sheet": name, "row_from": first_row, "row_to": last_row},
                locator_label=f'leht "{name}", read {first_row}–{last_row}',
            )
        )
        return total, produced + 1


def _cell_text(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "jah" if value else "ei"
    return normalise(str(value))


class PptxParser:
    name = "pptx"
    version = "1"
    mime_types = frozenset(
        {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
    )

    def parse(self, source: SourceFile) -> ParseResult:
        import pptx

        limits = current_limits()
        handle = io.BytesIO(source.content)
        guard_office_container(handle, limits)

        try:
            presentation = pptx.Presentation(handle)
        except Exception as error:
            raise ExtractionFailed(
                "unreadable_pptx", "PPTX-i ei õnnestunud avada; fail võib olla rikutud."
            ) from error

        fragments: list[Fragment] = []
        total = 0
        notes_found = 0

        for number, slide in enumerate(presentation.slides, start=1):
            visible = normalise(
                "\n".join(
                    shape.text_frame.text
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False)
                )
            )
            if visible:
                total += len(visible)
                guard_character_budget(total, limits)
                fragments.append(
                    Fragment(
                        text=visible,
                        locator_kind=LocatorKind.SLIDE,
                        locator={"slide": number},
                        locator_label=f"slaid {number}",
                    )
                )

            # Speaker notes are indexed, and labelled. They are frequently where
            # the actual position is written down — and they are also what the
            # presenter never intended the audience to read, so a result that
            # quotes them has to say where it got them.
            notes = _slide_notes(slide)
            if notes:
                notes_found += 1
                total += len(notes)
                guard_character_budget(total, limits)
                fragments.append(
                    Fragment(
                        text=notes,
                        locator_kind=LocatorKind.SLIDE,
                        locator={"slide": number, "part": "notes"},
                        locator_label=f"slaid {number}, esineja märkmed",
                    )
                )

        if not fragments:
            raise ExtractionFailed("no_text", "Esitlusest ei leitud teksti.")

        return ParseResult(
            derivatives=(
                DerivativePayload(
                    fragments=tuple(fragments),
                    metadata={
                        "slide_count": len(presentation.slides._sldIdLst),
                        "slides_with_notes": notes_found,
                    },
                ),
            )
        )


def _slide_notes(slide: object) -> str:
    if not getattr(slide, "has_notes_slide", False):
        return ""
    try:
        return normalise(slide.notes_slide.notes_text_frame.text)  # type: ignore[attr-defined]
    except Exception:
        return ""


registry.register(DocxParser())
registry.register(XlsxParser())
registry.register(PptxParser())
